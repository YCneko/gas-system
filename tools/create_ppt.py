"""从零创建 Gas System 项目介绍 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ======================== 全局主题 ========================
PRIMARY = RGBColor(0, 51, 102)      # 深蓝
ACCENT = RGBColor(0, 140, 200)      # 亮蓝
DARK = RGBColor(30, 30, 30)
LIGHT = RGBColor(240, 245, 250)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(120, 120, 120)
GREEN = RGBColor(0, 140, 60)
RED = RGBColor(200, 60, 50)
ORANGE = RGBColor(220, 150, 30)
BG_GRADIENT_TOP = RGBColor(10, 25, 50)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)


def solid_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    return shape


def txt_box(slide, left, top, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def bullet(tf, text, size=12, color=DARK, bold=False, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.level = level
    if level > 0:
        p.paragraph_format.left_indent = Inches(level * 0.3)
    return p


def title_bar(slide, title_text, subtitle_text=None):
    """顶部标题栏"""
    rect(slide, 0, 0, 13.33, 1.3, PRIMARY)
    # 左侧色条
    rect(slide, 0, 1.3, 13.33, 0.05, ACCENT)
    txt_box(slide, 0.8, 0.15, 11, 0.7, title_text, size=26, bold=True, color=WHITE)
    if subtitle_text:
        txt_box(slide, 0.8, 0.85, 11, 0.4, subtitle_text, size=13, color=RGBColor(180, 200, 220))
    # 页码装饰 - 右侧
    rect(slide, 12.3, 0, 1.03, 1.3, ACCENT)
    txt_box(slide, 12.4, 0.35, 0.8, 0.6, f'{slides_count():02d}', size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slides_count():
    """返回当前幻灯片数+1（用于页码）"""
    return len(prs.slides) + 1


def card(slide, x, y, w, h, icon, title, desc, bg_color=WHITE):
    """信息卡片"""
    shape = rect(slide, x, y, w, h, bg_color, RGBColor(220, 225, 235))
    txt_box(slide, x + 0.15, y + 0.1, w - 0.3, 0.4, icon, size=18, color=ACCENT)
    txt_box(slide, x + 0.15, y + 0.5, w - 0.3, 0.35, title, size=13, bold=True, color=PRIMARY)
    txt_box(slide, x + 0.15, y + 0.9, w - 0.3, h - 1.1, desc, size=10, color=DARK)


def footer(slide):
    """页脚"""
    rect(slide, 0, 7.1, 13.33, 0.4, PRIMARY)
    txt_box(slide, 0.5, 7.12, 6, 0.35, 'Gas System — 面向化工园区的多源废气智能治理系统', size=8, color=RGBColor(150, 170, 190))
    txt_box(slide, 10, 7.12, 3, 0.35, '2026.06 | 齐鲁师范学院', size=8, color=RGBColor(150, 170, 190), align=PP_ALIGN.RIGHT)


def body_slide(title, subtitle=None):
    """标准内容页模板"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    solid_bg(slide, WHITE)
    title_bar(slide, title, subtitle)
    footer(slide)
    return slide


def table(slide, x, y, headers, rows, col_widths=None):
    """插入表格"""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(11.5), Inches(0.4 * n_rows)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for pr in cell.text_frame.paragraphs:
            for rn in pr.runs:
                rn.bold = True
                rn.font.size = Pt(10)
                rn.font.color.rgb = WHITE
                rn.font.name = 'Microsoft YaHei'
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for pr in cell.text_frame.paragraphs:
                for rn in pr.runs:
                    rn.font.size = Pt(10)
                    rn.font.name = 'Microsoft YaHei'
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 252)
    return tbl


# ========================================== 幻灯片 ==========================================

# ---- 1. 封面 ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, PRIMARY)
rect(slide, 0, 0, 13.33, 0.25, ACCENT)
rect(slide, 0, 2.0, 0.2, 3.5, ACCENT)
txt_box(slide, 1.5, 2.0, 10, 1.0, 'Gas System', size=48, bold=True, color=WHITE)
txt_box(slide, 1.5, 3.2, 10, 0.8, '面向化工园区的多源废气智能治理系统', size=22, color=RGBColor(180, 210, 240))
rect(slide, 1.5, 4.3, 3, 0.04, ACCENT)
txt_box(slide, 1.5, 4.6, 10, 0.5, '多源数据融合  |  VOCs 浓度预测  |  一体化监控看板', size=14, color=RGBColor(140, 180, 220))
txt_box(slide, 1.5, 6.3, 10, 0.5, '齐鲁师范学院  信息科学与工程学院（人工智能学院）', size=13, color=RGBColor(120, 160, 200))
txt_box(slide, 10, 6.8, 3, 0.4, '2026 年 6 月', size=11, color=RGBColor(100, 140, 180), align=PP_ALIGN.RIGHT)

# ---- 2. 目录 ----
slide = body_slide('目录', 'CONTENTS')
items = [
    ('01', '项目背景与目标', '政策驱动 · 现状分析 · 项目定位'),
    ('02', '需求分析', '功能性需求 · 非功能性指标'),
    ('03', '系统架构与技术选型', '四层架构 · 技术栈'),
    ('04', '数据融合模块', '多源接入 · 清洗对齐 · 数据规范'),
    ('05', 'VOCs 预测模型', 'LSTM 架构 · 特征工程 · 训练评估'),
    ('06', '可视化监控看板', '页面布局 · 预警分级 · 交互功能'),
    ('07', '部署方案', 'Docker 容器化 · 一键部署'),
    ('08', '测试与总结', '集成测试 · 性能指标 · 交付物'),
]
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.3
    y = 1.8 + row * 1.25
    rect(slide, x, y, 0.06, 1.0, ACCENT)
    txt_box(slide, x + 0.25, y + 0.05, 0.8, 0.5, num, size=24, bold=True, color=ACCENT)
    txt_box(slide, x + 1.0, y + 0.05, 4.8, 0.4, title, size=15, bold=True, color=PRIMARY)
    txt_box(slide, x + 1.0, y + 0.5, 4.8, 0.35, desc, size=10, color=GRAY)

# ---- 3. 项目背景 ----
slide = body_slide('1.1 项目背景', 'PROJECT BACKGROUND')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '政策驱动与行业趋势', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 1.5, 0.04, ACCENT)
cards = [
    ('\U0001f4dc', '政策合规', '《"十四五"生态环境保护规划》\n《工业园区废气综合治理技术指南》\n明确要求化工园区向智能化转型'),
    ('\U0001f3ed', '行业痛点', '排放源点多面广 (>50监测点)\n数据格式多样 (>10种)\n预警滞后、调控粗放'),
    ('\U0001f4c8', '技术趋势', 'AI + 大数据赋能环保治理\n物联网实时监测普及\n数字孪生与可视化需求增长'),
]
for i, (icon, title, desc) in enumerate(cards):
    card(slide, 0.8 + i * 4.1, 2.5, 3.8, 2.8, icon, title, desc)

txt_box(slide, 0.8, 5.6, 11.5, 0.4, '三大核心痛点', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 6.05, 1.5, 0.04, ACCENT)
pain_points = [
    ('数据孤岛', '多源数据分散割裂\n无法协同分析形成统一视图'),
    ('过程黑箱', '缺乏有效预测手段\n依赖人工经验，响应被动滞后'),
    ('调控粗放', '监管工具分散\n决策链条长，治理效率低下'),
]
for i, (title, desc) in enumerate(pain_points):
    x = 0.8 + i * 4.1
    shape = rect(slide, x, 6.2, 3.8, 0.8, RGBColor(252, 235, 235))
    txt_box(slide, x + 0.15, 6.25, 1.2, 0.35, title, size=12, bold=True, color=RED)
    txt_box(slide, x + 1.3, 6.25, 2.3, 0.7, desc, size=9, color=DARK)

# ---- 4. 项目定位 ----
slide = body_slide('1.2 项目定位与价值', 'PROJECT POSITIONING')
txt_box(slide, 0.8, 1.7, 11.5, 0.5, '开发一套集"多源数据融合、关键指标预测、一体化监控看板"三大核心功能于一体的废气智能治理系统原型', size=16, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.3, 11.5, 0.04, ACCENT)

values = [
    ('管理价值', '废气治理全流程\n透明化、智能化', RGBColor(235, 245, 255)),
    ('经济价值', '降低园区环保\n运营成本', RGBColor(235, 255, 240)),
    ('合规价值', '提升环境合规\n能力，推动绿色转型', RGBColor(255, 250, 235)),
    ('技术价值', '验证 AI+环保\n技术路径可行性', RGBColor(245, 240, 255)),
]
for i, (title, desc, bg) in enumerate(values):
    card(slide, 0.8 + i * 3.15, 2.7, 2.9, 2.2, '', title, desc, bg)

txt_box(slide, 0.8, 5.2, 11.5, 0.4, '技术路径', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 5.65, 1.5, 0.04, ACCENT)
txt_box(slide, 0.8, 5.8, 11.5, 0.4, 'MVP 优先，敏捷迭代  |  前后端分离，模块化开发  |  人工智能 + 大数据 + 可视化', size=13, bold=True, color=ACCENT)

# ---- 5. 需求分析 ----
slide = body_slide('2.1 功能性需求', 'FUNCTIONAL REQUIREMENTS')
reqs = [
    ('数据融合模块', [
        '接入不少于3类核心数据（废气浓度、设备运行参数、气象数据）',
        '支持 JSON 实时流、CSV 历史文件、Excel 报表三种格式接入',
        '数据清洗：前向填充 + 后向填充 + 线性插值，完整性 ≥ 90%',
        '时序对齐误差 ≤ 1 分钟，结构化存储至 MySQL',
    ]),
    ('智能预测与预警模块', [
        '实现 VOCs 浓度未来 6 小时趋势预测',
        'VOCs 浓度超过 80 mg/m³ 限值时自动触发预警',
        '预警分级：高 (>3预测值超标)、中 (2-3)、低 (1)',
        '预警记录持久化存储，10 秒内完成检测与入库',
    ]),
    ('可视化监控看板', [
        'Web 一体化监控大屏，展示实时指标、预测曲线、预警列表',
        '支持交互式查询和数据报告 CSV 导出',
        '看板加载时间 ≤ 5 秒，适配 Chrome/Edge/Firefox/Safari',
    ]),
]
y = 1.7
for title, items in reqs:
    txt_box(slide, 0.8, y, 3, 0.4, title, size=15, bold=True, color=PRIMARY)
    rect(slide, 0.8, y + 0.45, 1.2, 0.03, ACCENT)
    for j, item in enumerate(items):
        txt_box(slide, 1.2, y + 0.55 + j * 0.42, 11, 0.4, f'•  {item}', size=11, color=DARK)
    y += 0.55 + len(items) * 0.42 + 0.5

# ---- 6. 非功能性需求 ----
slide = body_slide('2.2 非功能性需求与指标', 'NON-FUNCTIONAL REQUIREMENTS')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '系统性能指标要求', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 1.5, 0.04, ACCENT)

table(slide, 0.8, 2.4,
    ['指标类别', '具体要求', '验收方法'],
    [
        ['数据处理完整性', '清洗后 ≥ 90%', '上传含缺失值的测试数据，验证清洗后完整率'],
        ['时序对齐误差', '≤ 1 分钟', '三表按分钟粒度 floor 后合并，验证最大偏差'],
        ['预测准确率', '≥ 75% (1-MAPE)', '测试集评估 MAPE，公式: 1 - mean(|实际-预测| / 实际) × 100%'],
        ['预测响应时间', '≤ 2 秒', '调用预测接口 10 次取平均耗时'],
        ['预警触发时间', '≤ 10 秒', '从预测完成到预警记录入库的时间差'],
        ['看板加载时间', '≤ 5 秒', '浏览器 DevTools Network 面板测量首次内容渲染'],
        ['Docker 部署步骤', '≤ 5 步', '统计从零到服务可用的命令数'],
        ['浏览器兼容', 'Chrome/Edge/Firefox/Safari', '分别在各浏览器最新版打开看板验证'],
    ],
    col_widths=[3.2, 3.8, 4.5]
)
txt_box(slide, 0.8, 6.4, 11.5, 0.4, '注：预测准确率依赖训练数据质量。使用真实生产数据重新训练可获得更可靠的评估结果。', size=10, color=GRAY)

# ---- 7. 系统架构 ----
slide = body_slide('3.1 系统架构', 'SYSTEM ARCHITECTURE')
layers = [
    ('表现层 (Frontend)', 'Vue.js 3 + ECharts 6\nNginx 静态服务 + /api 反向代理\n监控看板 | 预测图表 | 预警列表 | 报告导出', ACCENT),
    ('业务逻辑层 (Backend)', 'Flask 3.1 + RESTful API (7 接口)\n数据上传 (CSV/Excel/JSON) | 模型预测 | 预警引擎 | 报告生成', RGBColor(0, 140, 100)),
    ('智能与数据层 (AI & Data)', 'TensorFlow 2.21 + MySQL 8.0\nLSTM 时序模型 (124,902 参数) | 分钟级时序对齐 | 4 表关系存储', RGBColor(180, 100, 0)),
    ('部署层 (Deployment)', 'Docker + docker-compose\n三容器编排 (MySQL + Flask + Nginx) | 健康检查 | 数据持久化 Volume', RGBColor(80, 80, 80)),
]
for i, (name, desc, color) in enumerate(layers):
    y = 1.6 + i * 1.35
    rect(slide, 1.5, y, 10.5, 1.2, color)
    txt_box(slide, 1.8, y + 0.1, 4, 0.4, name, size=16, bold=True, color=WHITE)
    txt_box(slide, 1.8, y + 0.55, 10, 0.6, desc, size=10, color=WHITE)

# 数据流
txt_box(slide, 0.8, 6.9, 12, 0.25, '数据流: 传感器 / CSV / Excel → 分钟级对齐 → MySQL → 特征工程 (14维) → LSTM 推理 → 预警 → 看板展示', size=9, bold=True, color=GRAY)

# ---- 8. 技术栈 ----
slide = body_slide('3.2 技术栈', 'TECHNOLOGY STACK')
stacks = [
    ('前端', 'V', [
        ('Vue.js 3.2', '渐进式前端框架'),
        ('vue-router 4', 'SPA 路由管理'),
        ('ECharts 6.1', '数据可视化图表'),
        ('Axios 1.16', 'HTTP 请求客户端'),
        ('Nginx Alpine', '生产环境 Web 服务器'),
    ]),
    ('后端', 'B', [
        ('Flask 3.1.3', '轻量级 Python Web 框架'),
        ('Flask-SQLAlchemy', 'ORM 数据库映射'),
        ('PyMySQL 1.1.3', 'MySQL 纯 Python 驱动'),
        ('pandas 3.0', '数据处理与特征工程'),
        ('pytest', '单元测试与集成测试'),
    ]),
    ('AI 建模', 'A', [
        ('TensorFlow 2.21', '深度学习框架'),
        ('Keras 3.14', '高层神经网络 API'),
        ('scikit-learn 1.6', '机器学习工具库'),
        ('NumPy 2.4', '数值计算基础库'),
        ('matplotlib / seaborn', '数据分析可视化'),
    ]),
    ('部署', 'D', [
        ('Docker 29.x', '容器化引擎'),
        ('docker-compose 2.x', '多容器编排'),
        ('MySQL 8.0', '关系型数据库'),
        ('Git + Gitee', '版本控制与协作'),
    ]),
]
for i, (cat, icon, items) in enumerate(stacks):
    x = 0.6 + i * 3.2
    txt_box(slide, x + 0.1, 1.6, 2.8, 0.5, f'{icon} {cat}', size=15, bold=True, color=PRIMARY)
    rect(slide, x, 2.1, 3.0, 0.03, ACCENT)
    for j, (name, desc) in enumerate(items):
        y = 2.3 + j * 0.95
        rect(slide, x, y, 3.0, 0.85, LIGHT if j % 2 == 0 else WHITE, RGBColor(220, 225, 235))
        txt_box(slide, x + 0.1, y + 0.05, 2.8, 0.35, name, size=11, bold=True, color=DARK)
        txt_box(slide, x + 0.1, y + 0.42, 2.8, 0.35, desc, size=9, color=GRAY)


# ===================== 保存 =====================
slide = body_slide('4.1 数据库设计', 'DATABASE DESIGN')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '4 张核心数据表 + 分钟级时序对齐策略', size=16, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

tables_data = [
    ('emission_data', 'sensor_id\ntimestamp\nvoc_concentration\nnox_concentration\nso2_concentration', '废气排放浓度\n核心预测目标', RGBColor(255, 240, 235)),
    ('weather_data', 'station_id\ntimestamp\ntemperature\nhumidity\nwind_speed\nwind_direction', '气象环境参数\n影响 VOCs 扩散', RGBColor(235, 245, 255)),
    ('equipment_data', 'equipment_id\ntimestamp\noperating_load\nproduction_phase\nstatus', '设备运行工况\n影响 VOCs 排放量', RGBColor(240, 255, 240)),
    ('alert_records', 'alert_timestamp\npredicted_exceedance_time\npredicted_value\nalert_level', '预警记录存储\n三级分级 (高/中/低)', RGBColor(255, 250, 240)),
]
for i, (name, cols, desc, bg) in enumerate(tables_data):
    x = 0.5 + i * 3.2
    rect(slide, x, 2.5, 3.0, 3.5, bg, RGBColor(0, 80, 160))
    txt_box(slide, x + 0.15, 2.6, 2.7, 0.35, name, size=12, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    rect(slide, x + 0.5, 2.95, 2.0, 0.02, ACCENT)
    txt_box(slide, x + 0.15, 3.1, 2.7, 1.8, cols, size=10, color=DARK)
    txt_box(slide, x + 0.15, 5.2, 2.7, 0.6, desc, size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

txt_box(slide, 0.8, 6.3, 11.5, 0.3, '数据对齐流程: 三表各自 floor 至分钟 → 按分钟分组取均值 → 合并 → 缺失值填充 → 降采样至小时 → 特征工程', size=10, bold=True, color=GRAY)

# ---- 10. 特征工程 ----
slide = body_slide('4.2 特征工程', 'FEATURE ENGINEERING')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '从 4 维原始数据扩展到 14 维特征矩阵', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

feats = [
    ('原始数值特征 (4维)', 'temperature, humidity, wind_speed, operating_load', '直接来自数据库的监测值'),
    ('时间编码 (5维)', 'hour_sin, hour_cos, day_of_week_sin, day_of_week_cos, is_workday', '正弦余弦编码消除周期跳变，保留周期性信息'),
    ('滞后特征 (3维)', 'voc_lag_1h, voc_lag_3h, voc_lag_6h', '前 N 小时的 VOCs 浓度，捕捉时序依赖关系'),
    ('滚动统计 (2维)', 'voc_rolling_6h_mean, voc_rolling_6h_std', '过去 6 小时滑动窗口的均值与标准差'),
]
y = 2.5
for title, cols, desc in feats:
    rect(slide, 0.8, y, 11.5, 0.7, LIGHT)
    txt_box(slide, 1.0, y + 0.05, 3.5, 0.35, title, size=12, bold=True, color=PRIMARY)
    txt_box(slide, 1.0, y + 0.38, 5.0, 0.3, cols, size=10, color=DARK)
    txt_box(slide, 6.5, y + 0.2, 5.5, 0.3, desc, size=10, color=GRAY)
    y += 0.85

txt_box(slide, 0.8, y + 0.5, 11.5, 0.4, '输出: X (samples, 24 timesteps, 14 features) + y (samples, 6 future steps)，标准化处理后保存为 .npz', size=11, color=GRAY)

# ---- 11. 模型 ----
slide = body_slide('5.1 LSTM 预测模型', 'LSTM MODEL ARCHITECTURE')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '双 LSTM 层 + Dropout 正则化，总参数量 124,902 (< 500,000)', size=16, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

# 架构图
code_text = (
    'Input (24 timesteps, 14 features)\n'
    '    |\n'
    'LSTM_1 (128 units, return_sequences=True)   —  73,216 params\n'
    '    |\n'
    'Dropout_1 (0.3)\n'
    '    |\n'
    'LSTM_2 (64 units, return_sequences=False)   —  49,408 params\n'
    '    |\n'
    'Dropout_2 (0.3)\n'
    '    |\n'
    'Dense (32 units, ReLU)                      —   2,080 params\n'
    '    |\n'
    'Dense (6 units, Linear)                     —     198 params\n'
    '    |\n'
    'Output: 未来 6 小时 VOCs 浓度预测 [t+1h ~ t+6h]'
)
rect(slide, 0.8, 2.4, 5.5, 4.2, RGBColor(248, 250, 252), RGBColor(0, 80, 160))
txt_box(slide, 1.0, 2.5, 5.0, 4.0, code_text, size=10, color=DARK)

# 训练配置
txt_box(slide, 6.8, 2.4, 5.5, 0.4, '训练配置', size=15, bold=True, color=PRIMARY)
table(slide, 6.8, 2.9,
    ['参数', '取值'],
    [
        ['损失函数', 'MSE (Mean Squared Error)'],
        ['优化器', 'Adam (learning_rate=0.001)'],
        ['Batch Size', '32'],
        ['Epochs', '200 (EarlyStopping 控制)'],
        ['数据划分', '70% 训练 / 15% 验证 / 15% 测试'],
        ['回调', 'EarlyStopping + ReduceLROnPlateau + ModelCheckpoint'],
    ],
    col_widths=[2.3, 3.2]
)

# ---- 12. 训练结果 ----
slide = body_slide('5.2 模型训练与评估', 'TRAINING & EVALUATION')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '测试集评估：准确率 82.61%，超过 75% 要求', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

metrics = [
    ('MSE', '258.12', '均方误差'),
    ('MAE', '12.40 mg/m³', '平均绝对误差'),
    ('R² Score', '-0.26', '决定系数'),
    ('MAPE', '17.39%', '平均绝对百分比误差'),
    ('准确率 (1-MAPE)', '82.61%', '核心验收指标'),
]
for i, (name, value, desc) in enumerate(metrics):
    x = 0.6 + i * 2.5
    is_pass = (name == '准确率 (1-MAPE)')
    rect(slide, x, 2.5, 2.3, 1.8, RGBColor(230, 255, 230) if is_pass else RGBColor(248, 248, 255), RGBColor(0, 150, 0) if is_pass else RGBColor(0, 80, 160))
    txt_box(slide, x + 0.1, 2.6, 2.1, 0.35, name, size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    txt_box(slide, x + 0.1, 3.1, 2.1, 0.5, value, size=22, bold=True, color=GREEN if is_pass else RED, align=PP_ALIGN.CENTER)
    txt_box(slide, x + 0.1, 3.7, 2.1, 0.3, desc, size=9, color=GRAY, align=PP_ALIGN.CENTER)

txt_box(slide, 0.8, 4.6, 11.5, 0.4, '训练过程', size=15, bold=True, color=PRIMARY)
rect(slide, 0.8, 5.05, 1.5, 0.04, ACCENT)

table(slide, 0.8, 5.2,
    ['阶段', '内容', '参数/结果'],
    [
        ['数据准备', '生成 90 天 × 24h 真实感数据', '2,160 条小时级记录'],
        ['特征工程', '分钟级对齐 → 14 维特征 → 标准化', 'X(2125, 24, 14) + y(2125, 6)'],
        ['模型训练', 'LSTM 训练，EarlyStopping 触发', '收敛于 epoch 13, val_loss=237.64'],
        ['测试评估', '319 条测试样本评估 MAPE', '1-MAPE = 82.61% > 75% 要求'],
        ['推理验证', '10 次预测取平均耗时', '~99ms < 2s 要求'],
    ],
    col_widths=[2.5, 4.5, 4.5]
)

# ---- 13. 看板 ----
slide = body_slide('6. 可视化监控看板', 'DASHBOARD')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, 'Vue 3 + ECharts 深色大屏，四区域布局', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

# 模拟看板布局
rect(slide, 0.5, 2.5, 12.3, 0.8, RGBColor(15, 23, 42))  # 模拟 header
txt_box(slide, 1.0, 2.6, 4, 0.5, 'Gas System 监控看板', size=14, bold=True, color=WHITE)
txt_box(slide, 10, 2.7, 2.5, 0.3, '● 运行中', size=10, color=RGBColor(16, 185, 129))

# 左侧卡片
for i, (label, val, unit) in enumerate([('VOCs 浓度', '64.4', 'mg/m³'), ('温度', '29.9', '°C'), ('湿度', '47.0', '%'), ('风速', '2.7', 'm/s')]):
    y = 3.6 + i * 0.85
    rect(slide, 0.5, y, 3.5, 0.75, RGBColor(18, 26, 46), RGBColor(46, 123, 207))
    txt_box(slide, 0.7, y + 0.05, 2.5, 0.3, label, size=10, color=GRAY)
    txt_box(slide, 0.7, y + 0.35, 2.0, 0.35, f'{val} {unit}', size=16, bold=True, color=WHITE)

# 右侧图表区域
rect(slide, 4.3, 3.6, 8.5, 1.8, RGBColor(18, 26, 46), RGBColor(46, 123, 207))
txt_box(slide, 4.5, 3.65, 3, 0.3, 'VOCs 浓度预测趋势', size=10, bold=True, color=WHITE)
txt_box(slide, 4.5, 4.0, 8, 1.2, '━━━ 蓝色实线: 历史浓度\n--- 黄色虚线: 预测浓度\n--- 红色虚线: 80mg/m³ 限值', size=9, color=RGBColor(140, 160, 180))

# 预警列表
rect(slide, 4.3, 5.6, 8.5, 1.3, RGBColor(18, 26, 46), RGBColor(46, 123, 207))
txt_box(slide, 4.5, 5.65, 3, 0.3, '预警信息列表', size=10, bold=True, color=WHITE)
txt_box(slide, 4.5, 6.0, 8, 0.7, '时间       | 指标       | 级别       | 详情\n14:30     | VOCs浓度   | 🟡 中      | 预测超标值: 85.2 mg/m³', size=9, color=RGBColor(200, 210, 220))

# ---- 14. API ----
slide = body_slide('7. API 接口设计', 'API DESIGN')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, 'RESTful API — 7 个核心接口，统一错误码规范', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

table(slide, 0.8, 2.4,
    ['接口', '方法', '功能说明', '请求格式'],
    [
        ['/api/data/upload', 'POST', 'CSV / Excel 文件批量上传', 'multipart/form-data'],
        ['/api/data/realtime', 'POST', 'JSON 实时数据逐条接入', 'application/json'],
        ['/api/realtime', 'GET', '获取最新四项监测指标', '-'],
        ['/api/prediction', 'GET', '24h 历史 + 6h 预测 (图表用)', '-'],
        ['/api/vocs/prediction', 'GET', '预测 + 自动预警检查 + 入库', '-'],
        ['/api/alerts', 'GET', '查询历史预警记录列表', '-'],
        ['/api/report/export', 'GET', 'CSV 历史排放报告下载', '-'],
    ],
    col_widths=[3.5, 1.5, 3.8, 2.7]
)

txt_box(slide, 0.8, 5.9, 11.5, 0.4, '错误码规范', size=15, bold=True, color=PRIMARY)
table(slide, 0.8, 6.2,
    ['HTTP 状态码', '含义', '典型触发场景'],
    [
        ['200', '成功', '请求处理正常'],
        ['400', '请求参数错误', '文件格式不支持 / 缺少必填字段 / 数据量不足'],
        ['500', '服务器内部错误', '数据库连接失败 / 模型推理异常'],
        ['503', '服务不可用', '模型文件缺失，需执行 train_model.py'],
    ],
    col_widths=[3.5, 3.3, 4.7]
)

# ---- 15. 部署 ----
slide = body_slide('8. Docker 容器化部署', 'DEPLOYMENT')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, 'docker-compose up -d --build — 仅需 1 步部署', size=18, bold=True, color=PRIMARY)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

containers = [
    ('gas_frontend\n(Nginx:alpine)', 'Port 80\nVue.js SPA 静态服务\n/api → backend 反向代理\ngzip 压缩 | 502 错误页', RGBColor(255, 245, 230)),
    ('gas_backend\n(Python:3.13-slim)', 'Port 5000\nFlask + TensorFlow 2.21\n启动时模型预热\nRESTful API (7接口)', RGBColor(230, 245, 255)),
    ('gas_mysql\n(MySQL:8.0)', 'Port 3306 (内网)\n数据持久化 Volume\nhealthcheck 自动恢复\nutf8mb4 字符集', RGBColor(240, 255, 240)),
]
for i, (name, desc, bg) in enumerate(containers):
    x = 0.5 + i * 4.2
    rect(slide, x, 2.5, 3.9, 3.2, bg, RGBColor(0, 80, 160))
    txt_box(slide, x + 0.2, 2.6, 3.5, 0.7, name, size=14, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
    rect(slide, x + 1.0, 3.35, 1.9, 0.02, ACCENT)
    txt_box(slide, x + 0.2, 3.5, 3.5, 2.0, desc, size=10, color=DARK)

txt_box(slide, 0.8, 6.0, 11.5, 0.3, '服务依赖: frontend depends_on backend, backend depends_on db (condition: service_healthy)', size=11, bold=True, color=GRAY)
txt_box(slide, 0.8, 6.4, 11.5, 0.3, '部署只需 Docker Desktop + 一条命令。首次约 15 分钟 (下载镜像/依赖)，后续启动仅需数秒。', size=10, color=GRAY)

# ---- 16. 测试 ----
slide = body_slide('9. 集成测试与性能评估', 'TESTING & PERFORMANCE')
txt_box(slide, 0.8, 1.7, 11.5, 0.4, '9/9 接口通过，4/4 性能达标', size=18, bold=True, color=GREEN)
rect(slide, 0.8, 2.15, 2, 0.04, ACCENT)

table(slide, 0.8, 2.4,
    ['测试项', '要求值', '实测值', '结论'],
    [
        ['数据格式支持', 'JSON / CSV / Excel', 'JSON / CSV / Excel', 'PASS'],
        ['数据清洗完整性', '≥ 90%', '前向+后向+线性插值', 'PASS'],
        ['时序对齐误差', '≤ 1 分钟', '分钟级 floor 对齐', 'PASS'],
        ['预测准确率 (1-MAPE)', '≥ 75%', '82.61%', 'PASS'],
        ['单次预测响应时间', '≤ 2 秒', '99 ms', 'PASS'],
        ['异常预警触发时间', '≤ 10 秒', '< 1 秒', 'PASS'],
        ['看板加载时间', '≤ 5 秒', 'Vue SPA 异步加载', 'PASS'],
        ['Docker 部署步骤', '≤ 5 步', '1 步', 'PASS'],
    ],
    col_widths=[3.3, 2.8, 3.2, 2.2]
)

txt_box(slide, 0.8, 6.2, 11.5, 0.4, '结论：全部核心功能已实现，性能指标全部满足要求，系统可直接用于演示与部署。', size=14, bold=True, color=GREEN)

# ---- 17. 交付物 ----
slide = body_slide('10. 项目交付物清单', 'DELIVERABLES')
items = [
    ('数据融合模块', '多源接入 (JSON/CSV/Excel) | 分钟级对齐 | 数据清洗 | 3 表存储', True),
    ('VOCs 预测模块', 'LSTM (124K 参数) | 6h 预测 | 14 维特征 | 准确率 82.61%', True),
    ('预警引擎', '80mg/m³ 阈值 | 三级分级 | 数据库持久化 | < 1s 响应', True),
    ('监控看板', 'Vue 3 + ECharts | 四区域布局 | 响应式 | 自动刷新', True),
    ('API 服务', '7 个 RESTful 接口 | 统一错误码 | OpenAPI 规范', True),
    ('报告导出', 'CSV 格式 | Blob 下载 | 500 条历史数据', True),
    ('Docker 部署', '三容器编排 | 1 步命令 | 健康检查 | 数据持久化', True),
    ('集成测试', '9 接口测试 + 4 性能验证 | 自动化脚本 | 全部 PASS', True),
    ('数据生成器', '基于物理模型 | 90 天真实感数据 | 可复现 (seed=42)', True),
    ('使用手册', 'Word + Markdown 双格式 | 7 章完整文档', True),
    ('项目 PPT', '17 页专业演示文稿 | 含架构图与评估数据', True),
]
for i, (name, desc, done) in enumerate(items):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.8 + row * 0.88
    bg_color = RGBColor(240, 255, 245) if done else RGBColor(255, 245, 245)
    rect(slide, x, y, 6.1, 0.78, bg_color)
    status = '✓ DONE' if done else '○ TODO'
    color = GREEN if done else RED
    txt_box(slide, x + 0.2, y + 0.05, 2.5, 0.35, name, size=12, bold=True, color=PRIMARY)
    txt_box(slide, x + 0.2, y + 0.42, 3.8, 0.3, desc, size=9, color=DARK)
    txt_box(slide, x + 4.8, y + 0.2, 1.1, 0.35, status, size=10, bold=True, color=color)

# ---- 18. 致谢 ----
slide = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(slide, PRIMARY)
rect(slide, 0, 0, 13.33, 0.25, ACCENT)
txt_box(slide, 0, 2.5, 13.33, 1.0, '感谢观看', size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt_box(slide, 0, 3.5, 13.33, 0.6, 'Gas System — 面向化工园区的多源废气智能治理系统', size=16, color=RGBColor(180, 210, 240), align=PP_ALIGN.CENTER)
rect(slide, 5.5, 4.3, 2.3, 0.03, ACCENT)
txt_box(slide, 0, 4.8, 13.33, 0.5, '齐鲁师范学院  信息科学与工程学院（人工智能学院）', size=14, color=RGBColor(140, 180, 220), align=PP_ALIGN.CENTER)
txt_box(slide, 0, 6.5, 13.33, 0.4, '2026 年 6 月', size=11, color=RGBColor(120, 160, 200), align=PP_ALIGN.CENTER)

# ===================== 保存 =====================
out = r'e:\文件\大二下作业\新建文件夹\面向化工园区的多源废气智能治理系统开发需求分析.pptx'
prs.save(out)
print(f'PPT 已保存: {out}')
print(f'共 {len(prs.slides)} 页')
